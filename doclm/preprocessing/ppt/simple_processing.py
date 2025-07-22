"""
module implements ppt parser
"""
import logging
# from langchain.document_loaders import UnstructuredPowerPointLoader as PptReader

from pptx import Presentation
from .util import get_file_info_from_reader

log = logging.getLogger("doclogger")


# pylint: disable=E1101
def simple_ppt_processor(stream, **kwargs):
    """
    processes Microsoft ppt format
    :param stream: BytesIO object
    :return:
    """

    prs = Presentation(stream)
    meta_info = get_file_info_from_reader(prs)
    log.debug('Extracted meta from ppt file')
    meta_info["pages"] = len([1 for slide in prs.slides if slide.element.get("show", None) is None])

    return meta_info, extract_slides(prs)


def extract_slides(prs):
    for slide in prs.slides:
        if slide.element.get("show", None) is not None:
            continue
        page_content = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    page_content += "\n" + run.text

        yield page_content
